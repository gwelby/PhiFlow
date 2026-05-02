#!/usr/bin/env python3
"""
Convert markdown slide deck to PowerPoint (.pptx)
Usage: python markdown_to_pptx.py <input.md> <output.pptx>
"""

import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Color scheme (dark blue primary, gold accents)
DARK_BLUE = RGBColor(26, 35, 126)  # #1a237e
GOLD = RGBColor(255, 215, 0)  # #ffd700
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(64, 64, 64)

def parse_slides(markdown_text):
    """Parse markdown into slides"""
    slides = []
    current_slide = None
    in_speaker_notes = False
    
    lines = markdown_text.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]
        
        # Detect slide start (## SLIDE N: or major heading after separator)
        if re.match(r'^## SLIDE \d+:', line, re.IGNORECASE):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'title': '',
                'content': [],
                'speaker_notes': [],
                'tables': []
            }
            in_speaker_notes = False
            i += 1
            continue
        
        # Detect speaker notes section
        if line.strip() == '**Speaker Notes:**':
            in_speaker_notes = True
            i += 1
            continue
        
        # Handle separator (---)
        if line.strip() == '---':
            if current_slide and not in_speaker_notes:
                slides.append(current_slide)
                current_slide = None
            in_speaker_notes = False
            i += 1
            continue
        
        # Skip empty lines
        if not line.strip():
            i += 1
            continue
        
        # Skip title slide header
        if line.startswith('# PhiFlow Pitch Deck'):
            i += 1
            continue
        
        # Main title (# Title)
        if line.startswith('# ') and current_slide:
            current_slide['title'] = line[2:].strip()
            i += 1
            continue
        
        # Subtitle (### Subtitle)
        if line.startswith('### ') and current_slide:
            current_slide['subtitle'] = line[4:].strip()
            i += 1
            continue
        
        # Content
        if current_slide:
            if in_speaker_notes:
                if line.strip().startswith('- '):
                    current_slide['speaker_notes'].append(line.strip()[2:])
            else:
                # Check for table
                if '|' in line:
                    # Parse table rows
                    table_lines = []
                    while i < len(lines) and '|' in lines[i]:
                        table_lines.append(lines[i])
                        i += 1
                    if table_lines:
                        current_slide['tables'].append(parse_table(table_lines))
                    continue
                else:
                    current_slide['content'].append(line)
        
        i += 1
    
    # Add last slide
    if current_slide:
        slides.append(current_slide)
    
    return slides

def parse_table(table_lines):
    """Parse markdown table into structured data"""
    rows = []
    for line in table_lines:
        if '---' in line.replace('|', ''):  # Skip separator rows
            continue
        cells = [cell.strip() for cell in line.split('|') if cell.strip()]
        if cells:
            rows.append(cells)
    return rows

def create_presentation(slides, output_path):
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for idx, slide_data in enumerate(slides):
        # Use blank layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Add title
        if slide_data.get('title'):
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(12.333), Inches(1)
            )
            tf = title_box.text_frame
            tf.text = slide_data['title']
            p = tf.paragraphs[0]
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
            p.alignment = PP_ALIGN.LEFT
        
        # Add subtitle
        if slide_data.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.5)
            )
            tf = subtitle_box.text_frame
            tf.text = slide_data['subtitle']
            p = tf.paragraphs[0]
            p.font.size = Pt(24)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.LEFT
        
        # Add content
        content_top = Inches(2) if slide_data.get('subtitle') else Inches(1.5)
        if slide_data.get('content'):
            content_box = slide.shapes.add_textbox(
                Inches(0.5), content_top, Inches(12.333), Inches(4)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for j, line in enumerate(slide_data['content']):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Handle bold text
                line = line.replace('**', '')
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(12)
        
        # Add tables
        if slide_data.get('tables'):
            for table_data in slide_data['tables']:
                if len(table_data) > 1:
                    rows = len(table_data)
                    cols = len(table_data[0]) if table_data else 0
                    
                    if cols > 0:
                        table = slide.shapes.add_table(
                            rows, cols,
                            Inches(0.5), Inches(2.5),
                            Inches(12), Inches(3)
                        ).table
                        
                        for r_idx, row in enumerate(table_data):
                            for c_idx, cell_text in enumerate(row):
                                if c_idx < cols:
                                    cell = table.cell(r_idx, c_idx)
                                    cell.text = cell_text.replace('**', '')
                                    
                                    # Style header row
                                    if r_idx == 0:
                                        cell.fill.solid()
                                        cell.fill.fore_color.rgb = DARK_BLUE
                                        for paragraph in cell.text_frame.paragraphs:
                                            paragraph.font.color.rgb = WHITE
                                            paragraph.font.bold = True
                                    else:
                                        for paragraph in cell.text_frame.paragraphs:
                                            paragraph.font.color.rgb = DARK_GRAY
        
        # Add speaker notes
        if slide_data.get('speaker_notes'):
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = '\n'.join(slide_data['speaker_notes'])
    
    # Save presentation
    prs.save(output_path)
    print(f"Saved PowerPoint: {output_path}")

def main():
    if len(sys.argv) < 3:
        print("Usage: python markdown_to_pptx.py <input.md> <output.pptx>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    with open(input_file, 'r', encoding='utf-8') as f:
        markdown_text = f.read()
    
    slides = parse_slides(markdown_text)
    print(f"Parsed {len(slides)} slides")
    
    create_presentation(slides, output_file)
    print(f"Created PowerPoint presentation: {output_file}")

if __name__ == '__main__':
    main()
